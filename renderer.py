"""
renderer.py — Playwright + python-pptx rendering core

Flow per slide:
  1. Build full URL: base_url + "/" + template
  2. Playwright opens URL, waits for data-render-complete="true"
  3. Inject theme via window.setTheme(theme)
  4. Inject data via window.setXxxData(data)
  5. Wait for any D3/chart re-renders to settle
  6. Screenshot 1920×1080 → PNG bytes
  7. python-pptx adds slide with PNG as full-bleed image
"""

import asyncio
import io
import json
import logging
from typing import Optional

from playwright.async_api import async_playwright, Page
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Presentation dimensions: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Viewport: must match HTML slide dimensions
VIEWPORT = {"width": 1920, "height": 1080}

# Timeout for slide render (ms)
RENDER_TIMEOUT = 30_000

# Map section_id → window inject function name
INJECT_MAP = {
    "divider":               "setDividerData",
    "disease_intro":         "setDiseaseData",
    "prevalence":            "setPrevalenceData",
    "treatment_landscape":   "setTreatmentData",
    "guidelines":            "setGuidelineData",
    "unmet_needs":           "setContentData",
    "moa":                   "setMoaData",
    "pivotal_studies":       "setPivotalData",
    "swimmer_plot":          "setSwimmerData",
    "subgroup_analysis":     "setSubgroupData",
    "competitive_landscape": "setCompetitiveData",
    "market_access":         "setMarketAccessData",
    "market_access_timeline":"setMarketAccessData",
    "treatment_algorithm":   "setAlgorithmData",
    "differentiators":       "setDifferentiatorsData",
    "swot":                  "setSwotData",
    "scientific_narrative":  "setNarrativeData",
    "strategic_imperatives": "setImperativesData",
    "tactical_plan":         "setTacticsData",
    "kol_mapping":           "setKOLData",
    "insights":              "setInsightsData",
    "timeline":              "setTimelineData",
    "executive_summary":     "setExecSummaryData",
}

# Extra settle time (ms) for slides with D3/animations
SETTLE_EXTRA = {
    "swimmer_plot":          800,
    "subgroup_analysis":     600,
    "competitive_landscape": 600,
    "kol_mapping":           1200,   # world-atlas map fetch
    "guidelines":            400,
}


async def screenshot_slide(
    page: Page,
    url: str,
    section_id: str,
    data: dict,
    theme: str,
) -> bytes:
    """
    Navigate to a slide URL, inject data + theme, return PNG bytes.
    """
    await page.goto(url, wait_until="networkidle", timeout=RENDER_TIMEOUT)

    # Wait for slide to signal it's ready
    try:
        await page.wait_for_selector(
            "[data-render-complete='true']",
            timeout=RENDER_TIMEOUT,
        )
    except Exception:
        logger.warning(f"[{section_id}] data-render-complete timeout — proceeding anyway")

    # Apply theme
    await page.evaluate(
        """(theme) => {
            if (typeof window.setTheme === 'function') {
                window.setTheme(theme);
            } else {
                // Fallback: set body class directly
                document.body.className = 'theme-' + theme;
            }
        }""",
        theme,
    )

    # Inject slide data if function exists
    inject_fn = INJECT_MAP.get(section_id)
    if inject_fn and data:
        injected = await page.evaluate(
            """([fn, payload]) => {
                if (typeof window[fn] === 'function') {
                    try { window[fn](payload); return true; }
                    catch(e) { return 'error: ' + e.message; }
                }
                return 'not_found';
            }""",
            [inject_fn, data],
        )
        if injected is not True:
            logger.warning(f"[{section_id}] inject {inject_fn}() → {injected}")

    # Extra settle time for complex renders
    extra = SETTLE_EXTRA.get(section_id, 0)
    if extra:
        await asyncio.sleep(extra / 1000)
    else:
        await asyncio.sleep(0.2)  # minimum settle

    # Hide theme switcher before screenshot
    await page.evaluate("""
        const sw = document.getElementById('themeSwitcher');
        if (sw) sw.style.display = 'none';
        const pc = document.getElementById('previewControls');
        if (pc) pc.style.display = 'none';
    """)

    # Screenshot full page at 1920×1080
    png = await page.screenshot(
        full_page=False,
        clip={"x": 0, "y": 0, "width": 1920, "height": 1080},
    )
    return png


async def render_presentation(request) -> bytes:
    """
    Render all slides and assemble PPTX. Returns bytes.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Blank slide layout
    blank_layout = prs.slide_layouts[6]

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-setuid-sandbox",
                "--disable-dev-shm-usage",
                "--disable-gpu",
                "--font-render-hinting=none",
            ],
        )
        context = await browser.new_context(
            viewport=VIEWPORT,
            device_scale_factor=1,
            # Ensure fonts render consistently
            extra_http_headers={"Accept-Language": "en-US,en;q=0.9"},
        )
        page = await context.new_page()

        total = len(request.slides)
        for idx, slide_cfg in enumerate(request.slides):
            url = f"{request.base_url}/{slide_cfg.template}"
            logger.info(f"[{idx+1}/{total}] {slide_cfg.section_id} → {url}")

            try:
                png = await screenshot_slide(
                    page=page,
                    url=url,
                    section_id=slide_cfg.section_id,
                    data=slide_cfg.data or {},
                    theme=slide_cfg.theme or request.theme,
                )
            except Exception as e:
                logger.error(f"[{idx+1}] Failed to render {slide_cfg.section_id}: {e}")
                # Insert error placeholder slide
                png = _error_png(slide_cfg.section_id, str(e))

            # Add slide to PPTX
            slide = prs.slides.add_slide(blank_layout)
            _add_full_bleed_image(slide, png, prs)

        await browser.close()

    # Serialize to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _add_full_bleed_image(slide, png_bytes: bytes, prs: Presentation):
    """Add a PNG as full-bleed background image to a slide."""
    img_stream = io.BytesIO(png_bytes)
    pic = slide.shapes.add_picture(
        img_stream,
        left=Emu(0),
        top=Emu(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )
    # Move image to back (z-order)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def _error_png(section_id: str, error: str) -> bytes:
    """Generate a simple error placeholder PNG using Pillow."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new("RGB", (1920, 1080), color=(15, 23, 42))
        draw = ImageDraw.Draw(img)
        draw.rectangle([60, 60, 1860, 1020], outline=(239, 68, 68), width=3)
        draw.text((80, 100), f"⚠ Render Error: {section_id}", fill=(239, 68, 68))
        draw.text((80, 140), error[:200], fill=(148, 163, 184))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf.read()
    except ImportError:
        # Fallback: 1×1 white PNG
        return (
            b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01'
            b'\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00'
            b'\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18'
            b'\xd8N\x00\x00\x00\x00IEND\xaeB`\x82'
        )
